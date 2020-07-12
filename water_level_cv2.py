import os
import glob
import numpy as np
import matplotlib.pyplot as plt
import cv2

# %% define a function for image capture


def image_capture(vidcap_obj, time_list):
    vidcap = vidcap_obj
    for idx, content in enumerate(time_list):
        # just cue to 20 sec. position
        vidcap.set(cv2.CAP_PROP_POS_MSEC, int(content*1000))
        success, image = vidcap.read()
        if success:
            cv2.imwrite("img_folder/frame_{}sec.jpg".format(content), image)
            plt.imshow(image)
            print("image captured successfully")
            cv2.waitKey(0)


# %%
vidcap = cv2.VideoCapture(
    'YouTube.mp4')
vidcap.set(cv2.CAP_PROP_POS_MSEC, 20000)      # just cue to 20 sec. position
success, image = vidcap.read()
if success:
    # save frame as JPEG file
    cv2.imwrite("img_folder/frame20sec7.jpg", image)
    # cv2.imshow("20sec",image)# will block the window
    plt.imshow(image)
    print("image captured successfully")
    # for 0, you need to activate your image and press any key. for 1 the function will show a frame for 1ms only
    cv2.waitKey(0)


# %% capture images from video with some time intervals.

time_list = np.arange(0, 880, 10)
image_capture(vidcap, time_list)

cv2.destroyAllWindows()

# %% clear figures, delete adjacent repeated figures.
while(True):
    FileNames = glob.glob("img_folder/*.jpg")
    FileNames_sorted = sorted(FileNames, key=os.path.getctime)
    idx = 28
    files_to_remove = []
    for idx, content in enumerate(FileNames_sorted):
        if len(FileNames_sorted) - idx < 2:
            break
        image_original = cv2.imread("{}".format(FileNames_sorted[idx]))
        #image_original = cv2.imread("img_folder/frame20sec7.jpg")
        image_to_compare = cv2.imread("{}".format(FileNames_sorted[idx+1]))

        if image_original.shape == image_to_compare.shape:
            #print("the images have same size and channels")
            image_difference = cv2.subtract(image_original, image_to_compare)
            blue, green, red = cv2.split(image_difference)
            # if cv2.countNonZero(blue) == 0 and cv2.countNonZero(green) == 0 and cv2.countNonZero(red) == 0:
            # if np.sum(blue) < 20e5 and np.sum(green) < 20e5 and np.sum(red) < 20e5:# algorithm1
            all_channels = blue + green + red
            if all_channels.max() < 250:  # or mean values
                print("{0} and {1} are equal, delete the second image.".format(
                    FileNames_sorted[idx], FileNames_sorted[idx+1]))
                # os.remove(FileNames_sorted[1])
                files_to_remove.append(FileNames_sorted[idx+1])
        else:
            # os.remove(FileNames_sorted[1])
            files_to_remove.append(FileNames_sorted[idx+1])

    # %%
    if len(files_to_remove) < 1:
        break
        # pass
    else:
        for content in files_to_remove:
            os.remove(content)
            print(content)


# %% compile all images to a ppt file
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
prs.slide_width = Inches(17.72)
prs.slide_height = Inches(10)
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "TIM introduction"
subtitle.text = "Prepare by PFL"

# prepare slide layout for images
blank_slide_layout = prs.slide_layouts[6]

FileNames = glob.glob("img_folder/*.jpg")
FileNames_sorted = sorted(FileNames, key=os.path.getctime)

for content in FileNames_sorted:
    img_path = content
    pic_slide = prs.slides.add_slide(blank_slide_layout)
    left = top = Inches(0)
    pic = pic_slide.shapes.add_picture(img_path, left, top)

prs.save("TIM introduction 01.pptx")

# %%
